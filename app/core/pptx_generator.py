from io import BytesIO
import requests
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.table import Table
import logging

from app.schemas.presentation import (
    Presentation, TitleSlide, BulletSlide, BulletPoint,
    ImageSlide, TableSlide, SplitSlide, ContentSection
)

# Local path to broken image placeholder
from pathlib import Path
BROKEN_IMAGE_PATH = Path(__file__).parent.parent / "static" / "images" / "broken-image.png"

# Set up logging
logger = logging.getLogger(__name__)


def create_presentation(presentation_data: Presentation) -> BytesIO:
    """Generate a PowerPoint presentation based on the provided schema."""
    prs = PPTXPresentation()
    
    for slide in presentation_data.slides:
        if slide.type == "title":
            create_title_slide(prs, slide)
        elif slide.type == "bullet":
            create_bullet_slide(prs, slide)
        elif slide.type == "image":
            create_image_slide(prs, slide)
        elif slide.type == "table":
            create_table_slide(prs, slide)
        elif slide.type == "split":
            create_split_slide(prs, slide)
    
    # Save the presentation to a BytesIO object
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    
    return output


def create_title_slide(prs: PPTXPresentation, slide_data: TitleSlide):
    """Create a title slide with a title and optional subtitle."""
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set the title
    title = slide.shapes.title
    title.text = slide_data.title
    
    # Set the subtitle if provided
    if slide_data.subtitle:
        subtitle = slide.placeholders[1]  # Subtitle placeholder
        subtitle.text = slide_data.subtitle


def create_bullet_slide(prs: PPTXPresentation, slide_data: BulletSlide):
    """Create a bullet point slide with an optional title and multi-level bullet points."""
    slide_layout = prs.slide_layouts[1]  # Bullet point layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set the title if provided
    if slide_data.title:
        title = slide.shapes.title
        title.text = slide_data.title
    
    # Add bullet points
    shapes = slide.shapes
    body_shape = shapes.placeholders[1]  # Content placeholder
    text_frame = body_shape.text_frame
    
    # Clear any existing text
    text_frame.clear()
    
    # Add bullet points
    for point in slide_data.points:
        add_bullet_point(text_frame, point, 0)


def add_bullet_point(text_frame, point: BulletPoint, current_level: int):
    """Recursively add bullet points with proper indentation."""
    p = text_frame.add_paragraph()
    p.text = point.text
    p.level = current_level  # Use the current nesting level
    
    # Add child bullet points recursively
    for child in point.children:
        add_bullet_point(text_frame, child, current_level + 1)


def create_image_slide(prs: PPTXPresentation, slide_data: ImageSlide):
    """Create a slide with an image."""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set the title if provided
    if slide_data.title:
        # Add a title textbox
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)
        
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        
        p = title_frame.add_paragraph()
        p.text = slide_data.title
        p.font.size = Pt(32)
        p.font.bold = True
    
    error_message = None
    
    # Try to get the image from URL
    try:
        response = requests.get(slide_data.url, timeout=10)
        response.raise_for_status()
        img_stream = BytesIO(response.content)
        image_source = "original"
    except Exception as e:
        # Log the original error
        error_message = f"Failed to load image from {slide_data.url}: {str(e)}"
        logger.warning(error_message)
        
        # Use local broken image placeholder
        try:
            with open(BROKEN_IMAGE_PATH, "rb") as f:
                img_stream = BytesIO(f.read())
            image_source = "placeholder"
        except Exception as err:
            # If broken image also fails, show error message
            logger.error(f"Failed to load broken image placeholder: {str(err)}")
            
            left = Inches(1)
            top = Inches(3)
            width = Inches(8)
            height = Inches(2)
            
            err_shape = slide.shapes.add_textbox(left, top, width, height)
            err_frame = err_shape.text_frame
            
            p = err_frame.add_paragraph()
            p.text = "Image unavailable"
            
            # Add presenter notes with error details
            if not slide.has_notes_slide:
                slide.notes_slide
            notes_text_frame = slide.notes_slide.notes_text_frame
            notes_text_frame.text = f"Error loading image: {error_message}\nBroken image placeholder also failed: {str(err)}"
            
            return
    
    # Add the image to the slide - centered
    left = Inches(1)
    top = Inches(2) if slide_data.title else Inches(1)
    width = Inches(8)  # Maximum width
    
    try:
        pic = slide.shapes.add_picture(img_stream, left, top, width=width)
        
        # Add alt text if provided
        if slide_data.alt:
            pic.name = slide_data.alt
            
        # Add presenter notes with error details if we used the placeholder
        if error_message and image_source == "placeholder":
            if not slide.has_notes_slide:
                slide.notes_slide
            notes_text_frame = slide.notes_slide.notes_text_frame
            notes_text_frame.text = f"Error loading image: {error_message}\nUsing broken image placeholder instead."
            
    except Exception as e:
        # If adding image to slide fails, add error message
        final_error = f"Failed to add image to slide: {str(e)}"
        logger.error(final_error)
        
        left = Inches(1)
        top = Inches(3)
        width = Inches(8)
        height = Inches(2)
        
        err_shape = slide.shapes.add_textbox(left, top, width, height)
        err_frame = err_shape.text_frame
        
        p = err_frame.add_paragraph()
        p.text = "Image unavailable"
        
        # Add presenter notes with error details
        if not slide.has_notes_slide:
            slide.notes_slide
        notes_text_frame = slide.notes_slide.notes_text_frame
        combined_error = f"Original error: {error_message if error_message else 'N/A'}\nFinal error: {final_error}"
        notes_text_frame.text = combined_error


def create_table_slide(prs: PPTXPresentation, slide_data: TableSlide):
    """Create a slide with a table."""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set the title if provided
    if slide_data.title:
        # Add a title textbox
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)
        
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        
        p = title_frame.add_paragraph()
        p.text = slide_data.title
        p.font.size = Pt(32)
        p.font.bold = True
    
    # Calculate table dimensions
    rows = len(slide_data.rows) + 1  # +1 for header row
    cols = len(slide_data.headers)
    
    # Add the table
    left = Inches(1)
    top = Inches(2) if slide_data.title else Inches(1)
    width = Inches(8)
    height = Inches(0.4 * rows)  # Adjust height based on number of rows
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Add headers
    for i, header in enumerate(slide_data.headers):
        cell = table.cell(0, i)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
    
    # Add data rows
    for i, row in enumerate(slide_data.rows):
        for j, cell_text in enumerate(row):
            if j < cols:  # Ensure we don't exceed the column count
                cell = table.cell(i + 1, j)  # +1 to skip header row
                cell.text = cell_text


def create_split_slide(prs: PPTXPresentation, slide_data: SplitSlide):
    """Create a slide with split content (two sections side by side)."""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set the title if provided
    title_height = 0
    if slide_data.title:
        # Add a title textbox
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)
        title_height = height
        
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_shape.text_frame
        
        p = title_frame.add_paragraph()
        p.text = slide_data.title
        p.font.size = Pt(32)
        p.font.bold = True
    
    # Layout parameters
    top = Inches(1.5) if slide_data.title else Inches(0.5)
    height = Inches(5)
    
    # Left section
    left_section = slide_data.sections[0]
    left = Inches(0.5)
    width = Inches(4.5)
    
    add_section_content(slide, left_section, left, top, width, height)
    
    # Right section
    right_section = slide_data.sections[1]
    left = Inches(5.0)
    width = Inches(4.5)
    
    add_section_content(slide, right_section, left, top, width, height)


def add_section_content(slide, section: ContentSection, left, top, width, height):
    """Add content to a section of a split slide."""
    if section.type == "bullet":
        # Add bullet points
        shape = slide.shapes.add_textbox(left, top, width, height)
        text_frame = shape.text_frame
        
        for point in section.points:
            add_bullet_point(text_frame, point, 0)
            
    elif section.type == "image":
        error_message = None
        
        try:
            # First try to get the requested image
            response = requests.get(section.url, timeout=10)
            response.raise_for_status()
            img_stream = BytesIO(response.content)
            image_source = "original"
        except Exception as e:
            # Log the original error
            error_message = f"Failed to load image from {section.url}: {str(e)}"
            logger.warning(error_message)
            
            # Use local broken image placeholder
            try:
                with open(BROKEN_IMAGE_PATH, "rb") as f:
                    img_stream = BytesIO(f.read())
                image_source = "placeholder"
            except Exception as err:
                # If broken image also fails, show error message
                placeholder_error = f"Failed to load broken image placeholder: {str(err)}"
                logger.error(placeholder_error)
                
                shape = slide.shapes.add_textbox(left, top, width, height)
                text_frame = shape.text_frame
                
                p = text_frame.add_paragraph()
                p.text = "Image unavailable"
                
                # Add presenter notes with error details
                if not slide.has_notes_slide:
                    slide.notes_slide
                notes_text_frame = slide.notes_slide.notes_text_frame
                notes_text = notes_text_frame.text if notes_text_frame.text else ""
                notes_text += f"\nSection image error: {error_message}\nBroken image error: {placeholder_error}"
                notes_text_frame.text = notes_text.strip()
                
                return
        
        try:
            # Add the image
            pic = slide.shapes.add_picture(img_stream, left, top, width=width)
            
            # Add alt text if provided
            if section.alt:
                pic.name = section.alt
                
            # Add presenter notes with error details if we used the placeholder
            if error_message and image_source == "placeholder":
                if not slide.has_notes_slide:
                    slide.notes_slide
                notes_text_frame = slide.notes_slide.notes_text_frame
                notes_text = notes_text_frame.text if notes_text_frame.text else ""
                notes_text += f"\nSection image error: {error_message}"
                notes_text_frame.text = notes_text.strip()
                
        except Exception as e:
            # If adding image to slide fails, add error message
            final_error = f"Failed to add image to split slide: {str(e)}"
            logger.error(final_error)
            
            shape = slide.shapes.add_textbox(left, top, width, height)
            text_frame = shape.text_frame
            
            p = text_frame.add_paragraph()
            p.text = "Image unavailable"
            
            # Add presenter notes with error details
            if not slide.has_notes_slide:
                slide.notes_slide
            notes_text_frame = slide.notes_slide.notes_text_frame
            notes_text = notes_text_frame.text if notes_text_frame.text else ""
            notes_text += f"\nSection image original error: {error_message if error_message else 'N/A'}"
            notes_text += f"\nSection image final error: {final_error}"
            notes_text_frame.text = notes_text.strip()
            
    elif section.type == "table":
        if not section.headers or not section.rows:
            return
        
        # Calculate table dimensions
        rows = len(section.rows) + 1  # +1 for header row
        cols = len(section.headers)
        
        # Add the table
        table_height = min(height, Inches(0.4 * rows))
        table = slide.shapes.add_table(rows, cols, left, top, width, table_height).table
        
        # Add headers
        for i, header in enumerate(section.headers):
            cell = table.cell(0, i)
            cell.text = header
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
        
        # Add data rows
        for i, row in enumerate(section.rows):
            for j, cell_text in enumerate(row):
                if j < cols:  # Ensure we don't exceed the column count
                    cell = table.cell(i + 1, j)  # +1 to skip header row
                    cell.text = cell_text